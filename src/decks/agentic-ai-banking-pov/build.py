"""
Build the 23-slide Agentic AI in Banking POV deck.

Aesthetic: financial-research-paper meets top-tier consulting broadsheet.
Same visual vocabulary as the primer (Fraunces + Newsreader + JetBrains
Mono, the "seam" motif, hairline-heavy, one BLUE accent per slide).

Colors and the font naming policy are sourced from ../../tokens.js
(EDITORIAL family). DEEP_BLUE appears only in the thin title band at the
top of each content slide — never as a column header, lane label, cell
fill, or chart series.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─────────────────────────────────────────────────────────────────────
# IC tokens — mirrored from ../../tokens.js
# ─────────────────────────────────────────────────────────────────────
def rgb(h): return RGBColor.from_string(h.lstrip("#"))

T = {
    "WHITE":         rgb("FFFFFF"),
    "DEEP":          rgb("182534"),
    "PAVEMENT":      rgb("404040"),
    "BLUE":          rgb("007CC3"),
    "STONE":         rgb("A7A9AC"),
    "SLATE":         rgb("6D6E71"),
    "LT_BLUE":       rgb("3699DA"),
    "TINT_BLUE":     rgb("E6F2FA"),
    "TINT_STONE":    rgb("F2F2F3"),
    "RULE":          rgb("E2E4E7"),
    "RED":           rgb("C0392B"),
    "AMBER":         rgb("D89B1F"),
    "GREEN":         rgb("2E8B57"),
}

# Editorial font stack — installed locally, will fall back gracefully if
# the recipient doesn't have them.
FONT_DISPLAY = "Fraunces"
FONT_BODY    = "Newsreader"
FONT_MONO    = "JetBrains Mono"

# Canvas — 16:9, 13.333" × 7.5"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Generous margins (consulting decks breathe)
MARGIN   = Inches(0.7)
BAND_H   = Inches(0.5)        # thin title band — more elegant than the 0.667 default
BODY_TOP = BAND_H + Inches(0.55)
BODY_BOT = SLIDE_H - Inches(0.55)
BODY_W   = SLIDE_W - 2 * MARGIN
BODY_H   = BODY_BOT - BODY_TOP

# ─────────────────────────────────────────────────────────────────────
# Primitive helpers
# ─────────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, *, line=None, line_w=0.5):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill if isinstance(fill, RGBColor) else T[fill]
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line if isinstance(line, RGBColor) else T[line]
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s

def hline(slide, x1, y, x2, color="RULE", weight=0.5):
    ln = slide.shapes.add_connector(1, x1, y, x2, y)
    ln.line.color.rgb = color if isinstance(color, RGBColor) else T[color]
    ln.line.width = Pt(weight)
    return ln

def vline(slide, x, y1, y2, color="RULE", weight=0.5):
    ln = slide.shapes.add_connector(1, x, y1, x, y2)
    ln.line.color.rgb = color if isinstance(color, RGBColor) else T[color]
    ln.line.width = Pt(weight)
    return ln

def text(slide, x, y, w, h, body, *, font=FONT_BODY, size=12, bold=False,
         italic=False, color="PAVEMENT", align="left", anchor="top",
         line_spacing=1.35, letter_spacing=None, font_weight=None):
    """Add a text box. `body` is a string OR a list of (text, opts) tuples for
    inline runs in a single paragraph, OR a list of strings/lists for multiple
    paragraphs."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]

    # Normalize: paragraphs is a list of (text-or-runs) entries.
    if isinstance(body, str):
        paragraphs = [body]
    else:
        paragraphs = body

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                 "right": PP_ALIGN.RIGHT}

    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map[align]
        p.line_spacing = line_spacing

        # A paragraph can be: str | list of (txt, opts) | (txt, opts)
        if isinstance(para, str):
            runs = [(para, {})]
        elif isinstance(para, tuple):
            runs = [para]
        else:
            runs = para

        for t_str, o in runs:
            r = p.add_run()
            r.text = t_str
            r.font.name = o.get("font", font)
            r.font.size = Pt(o.get("size", size))
            r.font.bold = o.get("bold", bold)
            r.font.italic = o.get("italic", italic)
            c = o.get("color", color)
            r.font.color.rgb = c if isinstance(c, RGBColor) else T[c]
    return tb


def oval(slide, x, y, w, h, fill, *, line=None, line_w=0.5):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill if isinstance(fill, RGBColor) else T[fill]
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line if isinstance(line, RGBColor) else T[line]
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


# ─────────────────────────────────────────────────────────────────────
# Master frame helpers
# ─────────────────────────────────────────────────────────────────────
def title_band(slide, title, eyebrow=None):
    """Thin DEEP_BLUE title band. Fraunces title, mono eyebrow on the right.
    This is one of two allowed uses of DEEP_BLUE in the deck (the other is
    the cover footer strip)."""
    rect(slide, 0, 0, SLIDE_W, BAND_H, "DEEP")
    text(slide, MARGIN, 0, SLIDE_W - 2*MARGIN, BAND_H, title,
         font=FONT_DISPLAY, size=18, bold=False, italic=False,
         color="WHITE", anchor="middle", line_spacing=1.05)
    if eyebrow:
        text(slide, SLIDE_W - Inches(4) - MARGIN, 0, Inches(4), BAND_H,
             eyebrow.upper(), font=FONT_MONO, size=8, bold=True,
             color="STONE", align="right", anchor="middle",
             letter_spacing=0.18)


def footer(slide, n, total=25):
    """Hairline rule + mono page metadata."""
    y = SLIDE_H - Inches(0.38)
    hline(slide, MARGIN, y - Inches(0.05), SLIDE_W - MARGIN,
          color="RULE", weight=0.5)
    text(slide, MARGIN, y, Inches(8), Inches(0.3),
         "INFOSYS CONSULTING  ·  AGENTIC AI IN BANKING  ·  POINT OF VIEW",
         font=FONT_MONO, size=7, bold=False, color="STONE")
    text(slide, SLIDE_W - Inches(2) - MARGIN, y, Inches(2), Inches(0.3),
         f"§ {n:02d} / {total:02d}",
         font=FONT_MONO, size=7, bold=True, color="BLUE",
         align="right")


def placeholder_pill(slide, label):
    """Compact placeholder pill — top-right of body area, BLUE rule + tint."""
    w = Inches(3.4)
    h = Inches(0.26)
    x = SLIDE_W - MARGIN - w
    y = BAND_H + Inches(0.15)
    rect(slide, x, y, w, h, "TINT_BLUE")
    rect(slide, x, y, Inches(0.04), h, "BLUE")
    text(slide, x + Inches(0.12), y, w - Inches(0.15), h,
         f"PLACEHOLDER  ·  {label.upper()}",
         font=FONT_MONO, size=7, bold=True, color="BLUE",
         anchor="middle")


def chapter_eyebrow(slide, chapter, label):
    """Small mono chapter marker placed just under the title band, left side.
    Reinforces the editorial chapter structure across the deck. No rule below
    — the rule was visually conflicting with body text on several slides."""
    y = BAND_H + Inches(0.18)
    text(slide, MARGIN, y, Inches(8), Inches(0.2),
         f"CHAPTER  {chapter}   ·   {label.upper()}",
         font=FONT_MONO, size=8, bold=True, color="BLUE",
         letter_spacing=0.16)


# ─────────────────────────────────────────────────────────────────────
# Build the presentation
# ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title slide (journal-cover treatment)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# Top DEEP_BLUE masthead band — the publication's "frontispiece"
mast_h = Inches(0.55)
rect(s, 0, 0, SLIDE_W, mast_h, "DEEP")
text(s, MARGIN, 0, Inches(7), mast_h,
     "INFOSYS  CONSULTING",
     font=FONT_MONO, size=10, bold=True, color="WHITE",
     anchor="middle", letter_spacing=0.22)
text(s, SLIDE_W - MARGIN - Inches(7), 0, Inches(7), mast_h,
     "POINT OF VIEW  ·  NO. 01  ·  BANKING  ·  MAY 2026",
     font=FONT_MONO, size=9, color="STONE",
     align="right", anchor="middle", letter_spacing=0.16)

# Mono kicker tag — research-note style
text(s, MARGIN, Inches(1.55), BODY_W, Inches(0.3),
     "A POINT OF VIEW  ·  AGENTIC AI IN BANKING",
     font=FONT_MONO, size=10, bold=True, color="BLUE",
     letter_spacing=0.24)
# Short BLUE seam under kicker
rect(s, MARGIN, Inches(1.95), Inches(0.6), Inches(0.025), "BLUE")

# Hero title — Fraunces, large, italic BLUE emphasis on "work"
text(s, MARGIN, Inches(2.4), BODY_W, Inches(2.8),
     [[("Agents that change the ", {"color": "DEEP"}),
       ("work", {"italic": True, "color": "BLUE"}),
       (",", {"color": "DEEP"})],
      [("not just the worker.", {"color": "DEEP"})]],
     font=FONT_DISPLAY, size=58, bold=False, color="DEEP",
     line_spacing=1.08)

# Subtitle in italic Newsreader — clear of descenders
text(s, MARGIN, Inches(5.25), Inches(10), Inches(0.8),
     "A primer on agentic AI in banking — and why the next decade of "
     "value comes from redesigning processes, not buying better assistants.",
     font=FONT_BODY, size=15, italic=True, color="SLATE",
     line_spacing=1.5)

# Bottom byline grid — 4-column metadata in mono + Newsreader.
# Sits above a thin DEEP_BLUE colophon strip at the very bottom.
byline_y = Inches(6.3)
hline(s, MARGIN, byline_y, SLIDE_W - MARGIN, color="RULE", weight=0.5)
fields = [
    ("PRACTICE",     "Banking"),
    ("EDITION",      "No. 01  ·  May 2026"),
    ("AUDIENCE",     "CXO buyers  ·  IC pursuit teams"),
    ("DISTRIBUTION", "Sales enablement — Internal"),
]
col_w = (BODY_W) / 4
for i, (label, value) in enumerate(fields):
    fx = MARGIN + i * col_w
    text(s, fx, byline_y + Inches(0.18), col_w - Inches(0.2), Inches(0.2),
         label,
         font=FONT_MONO, size=8, bold=True, color="STONE",
         letter_spacing=0.18)
    text(s, fx, byline_y + Inches(0.42), col_w - Inches(0.2), Inches(0.4),
         value,
         font=FONT_BODY, size=12, italic=True, color="PAVEMENT",
         line_spacing=1.3)

# Thin DEEP_BLUE colophon at the very bottom — the cover's only other
# allowed DEEP_BLUE element.
strip_h = Inches(0.35)
rect(s, 0, SLIDE_H - strip_h, SLIDE_W, strip_h, "DEEP")
text(s, MARGIN, SLIDE_H - strip_h, BODY_W, strip_h,
     "INFOSYS CONSULTING  ·  BANKING PRACTICE  ·  AGENTIC AI IN BANKING",
     font=FONT_MONO, size=7, bold=True, color="STONE", anchor="middle",
     letter_spacing=0.2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — Contents (Table of contents)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "Contents", "—")
chapter_eyebrow(s, "—", "Table of contents")

# Subtitle
text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.5),
     "Six chapters and an appendix. About 25 minutes if read in sequence.",
     font=FONT_DISPLAY, size=20, italic=True, color="DEEP", line_spacing=1.1)

# TOC rows — Roman numeral, chapter title, descriptive blurb, page number
contents = [
    ("I",    "Setup",            "Forty pilots, value in two",                  "03"),
    ("II",   "Reframe",          "The unit of work has changed",                "06"),
    ("III",  "Proof",            "Three banking processes, redrawn",             "09"),
    ("IV",   "Maturity & Path",  "Where banks are, where they're going",         "13"),
    ("V",    "How IC Helps",     "Capabilities, conversation tools, examples",   "17"),
    ("VI",   "Close",            "The bottom line, and what to do next",         "22"),
    ("—",    "Appendix",         "Sources and references",                       "24"),
]
top = BODY_TOP + Inches(1.05)
row_h = Inches(0.62)
for i, (roman, title_text, sub, page) in enumerate(contents):
    y = top + i * row_h
    # Top hairline rule for each row
    hline(s, MARGIN, y, MARGIN + BODY_W, color="RULE", weight=0.5)
    # Roman numeral in Fraunces italic BLUE
    text(s, MARGIN, y + Inches(0.12), Inches(0.9), Inches(0.45), roman,
         font=FONT_DISPLAY, size=22, italic=True, color="BLUE",
         line_spacing=1.0)
    # Chapter title in Fraunces
    text(s, MARGIN + Inches(1.1), y + Inches(0.13), Inches(4.0), Inches(0.4),
         title_text,
         font=FONT_DISPLAY, size=19, color="DEEP", line_spacing=1.1)
    # Descriptive blurb in Newsreader italic
    text(s, MARGIN + Inches(5.2), y + Inches(0.18),
         BODY_W - Inches(6.0), Inches(0.4),
         sub,
         font=FONT_BODY, size=13, italic=True, color="SLATE",
         line_spacing=1.3)
    # Page number in mono BLUE, right-aligned
    text(s, SLIDE_W - MARGIN - Inches(0.8), y + Inches(0.16),
         Inches(0.8), Inches(0.4), page,
         font=FONT_MONO, size=14, bold=True, color="BLUE",
         align="right", line_spacing=1.0)
# Closing hairline
hline(s, MARGIN, top + len(contents)*row_h, MARGIN + BODY_W,
      color="RULE", weight=0.5)

footer(s, 2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Executive summary
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "Three things to take from this deck",
           "Executive summary")
chapter_eyebrow(s, "I", "Setup")

# Slide-level subtitle in Fraunces italic
text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.5),
     "The argument, in three lines.",
     font=FONT_DISPLAY, size=22, italic=True, color="DEEP",
     line_spacing=1.1)

# Three numbered cards — hairlines, not boxes
items = [
    ("01", "THE TRAP",
     "The median bank runs 40+ generative AI pilots; realized value "
     "concentrates in only two or three. The pilot factory has stopped "
     "paying off."),
    ("02", "THE REFRAME",
     "Pilots plateau because agents inherit the friction of the process "
     "they joined. The prize is redesigning the process around what an "
     "agent can actually do."),
    ("03", "THE PATH",
     "Five no-regret moves now, one redesigned process within twelve "
     "months, and an operating-model shift across the next twenty-four."),
]
top = BODY_TOP + Inches(1.0)
card_h = Inches(4.5)
card_w = (BODY_W - Inches(0.8)) / 3
for i, (num, head, body) in enumerate(items):
    x = MARGIN + i * (card_w + Inches(0.4))
    # Top thin BLUE rule (the seam)
    rect(s, x, top, Inches(0.5), Inches(0.025), "BLUE")
    # Numeric anchor in Fraunces
    text(s, x, top + Inches(0.25), card_w, Inches(1.4), num,
         font=FONT_DISPLAY, size=80, bold=False, italic=True,
         color="BLUE", line_spacing=1.0)
    # Mono eyebrow
    text(s, x, top + Inches(1.7), card_w, Inches(0.3), head,
         font=FONT_MONO, size=9, bold=True, color="DEEP",
         letter_spacing=0.18)
    # Body in Newsreader
    text(s, x, top + Inches(2.1), card_w, Inches(2.4), body,
         font=FONT_BODY, size=14, color="PAVEMENT", line_spacing=1.5)

footer(s, 3)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Why now (3 columns)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "Three things changed in the last 18 months",
           "Setup  ·  Why now")
chapter_eyebrow(s, "I", "Why now")

# Slide subtitle
text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.5),
     "A phase change, not a headline.",
     font=FONT_DISPLAY, size=22, italic=True, color="DEEP",
     line_spacing=1.1)

cols = [
    ("i.", "Capability",
     "Frontier models can plan, use tools, and verify their own work. "
     "Per-token cost has fallen ~85% since 2023, putting capable models "
     "in reach of bulk operational workloads.",
     "Anthropic, OpenAI, Google · 2024-2026"),
    ("ii.", "Tooling",
     "Orchestration, evaluation harnesses, observability, and guardrails "
     "have moved from research-interesting to production-engineering. "
     "Banks can buy or self-host an agent stack like they buy analytics.",
     "LangGraph · OpenAI Assistants · Anthropic"),
    ("iii.", "Maturity",
     "JPMorgan reports $17B+ in annual tech spend with AI as a named "
     "line item. BofA's Erica has logged 2B+ client interactions. The "
     "infrastructure is in place; the design discipline is not.",
     "JPM 10-K 2024 · BCG Build for the Future 2024"),
]
top = BODY_TOP + Inches(1.0)
col_w = (BODY_W - Inches(0.8)) / 3
for i, (roman, head, body, src) in enumerate(cols):
    x = MARGIN + i * (col_w + Inches(0.4))
    # Thin BLUE rule top
    rect(s, x, top, Inches(0.4), Inches(0.025), "BLUE")
    # Roman numeral
    text(s, x, top + Inches(0.2), col_w, Inches(0.6), roman,
         font=FONT_DISPLAY, size=32, italic=True, color="BLUE",
         line_spacing=1.0)
    # Fraunces head
    text(s, x, top + Inches(0.95), col_w, Inches(0.7), head,
         font=FONT_DISPLAY, size=24, color="DEEP", line_spacing=1.1)
    # Newsreader body
    text(s, x, top + Inches(1.85), col_w - Inches(0.2), Inches(2.6),
         body, font=FONT_BODY, size=13, color="PAVEMENT", line_spacing=1.55)
    # Mono source
    text(s, x, BODY_BOT - Inches(0.6), col_w - Inches(0.2), Inches(0.4),
         src.upper(), font=FONT_MONO, size=7, color="STONE",
         letter_spacing=0.14)

footer(s, 4)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — The trap (long-tail chart as research figure)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "Forty pilots, value in two", "Setup  ·  The trap")
chapter_eyebrow(s, "I", "The trap")

# Slide subtitle — large Fraunces with BLUE italic emphasis
text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W - Inches(2.5), Inches(1.2),
     [[("The pilot factory has ", {"color": "DEEP"}),
       ("stopped", {"italic": True, "color": "BLUE"}),
       (" paying off.", {"color": "DEEP"})]],
     font=FONT_DISPLAY, size=30, color="DEEP", line_spacing=1.1)

# Right-side annotation
text(s, SLIDE_W - MARGIN - Inches(3.5), BODY_TOP + Inches(0.15),
     Inches(3.5), Inches(1.0),
     "Across the largest banks, a long-tail distribution: a small "
     "number of use cases deliver concentrated value while most pilots "
     "remain stuck in proof-of-concept.",
     font=FONT_BODY, size=11, italic=True, color="SLATE", line_spacing=1.5)

# Research-paper-style figure box
fig_top = BODY_TOP + Inches(1.7)
fig_h = Inches(3.8)
fig_x = MARGIN
fig_w = BODY_W

# Draw the figure box FIRST so it doesn't cover the FIG.01 label
rect(s, fig_x, fig_top, fig_w, fig_h, "WHITE", line="RULE", line_w=0.5)
# FIG label clearly above the box top border
text(s, fig_x + Inches(0.3), fig_top - Inches(0.32), Inches(2), Inches(0.25),
     "FIG.  01", font=FONT_MONO, size=8, bold=True, color="BLUE",
     letter_spacing=0.2)

# Figure title (Fraunces)
text(s, fig_x + Inches(0.3), fig_top + Inches(0.15),
     fig_w - Inches(0.6), Inches(0.4),
     "Pilot value concentrates in a long tail.",
     font=FONT_DISPLAY, size=14, color="DEEP", line_spacing=1.2)
text(s, fig_x + Inches(0.3), fig_top + Inches(0.55),
     fig_w - Inches(0.6), Inches(0.35),
     "Per-bank distribution of realized value across active "
     "generative-AI initiatives. The top two carry the majority.",
     font=FONT_BODY, size=10, italic=True, color="SLATE", line_spacing=1.4)

# Chart bars
chart_x = fig_x + Inches(0.5)
chart_y = fig_top + Inches(1.15)
chart_w = fig_w - Inches(1.0)
chart_h = Inches(2.3)

# Baseline
hline(s, chart_x, chart_y + chart_h, chart_x + chart_w,
      color="PAVEMENT", weight=0.75)

# Long-tail bar heights (relative to chart_h)
heights = [1.00, 0.84, 0.46, 0.27, 0.19, 0.15, 0.13, 0.115, 0.10, 0.09,
           0.082, 0.075, 0.068, 0.062, 0.057, 0.052, 0.048, 0.044, 0.04,
           0.037, 0.034, 0.031, 0.028, 0.025]
n = len(heights)
gap = Emu(40000)
total_gap = gap * (n - 1)
bar_w = Emu((chart_w - total_gap) // n)
for i, h_rel in enumerate(heights):
    bh = Emu(int(chart_h * h_rel))
    bx = chart_x + i * (bar_w + gap)
    by = chart_y + chart_h - bh
    if i < 2:
        c = "BLUE"
    elif i == 2:
        c = "LT_BLUE"
    else:
        c = "STONE"
    rect(s, bx, by, bar_w, bh, c)

# Top-left callout — placed clearly ABOVE the tallest bar
text(s, chart_x, chart_y - Inches(0.4), Inches(5.0), Inches(0.25),
     "↘  ~2 USE CASES CARRY THE VALUE",
     font=FONT_MONO, size=9, bold=True, color="BLUE", letter_spacing=0.16)

# Tail annotation — placed lower-right where bars are tiny
text(s, chart_x + Inches(4.5), chart_y + Inches(0.6),
     chart_w - Inches(4.5), Inches(0.5),
     "the long tail — busy, expensive, ineffective",
     font=FONT_BODY, size=11, italic=True, color="SLATE")

# Axis labels (mono)
text(s, chart_x, chart_y + chart_h + Inches(0.15), Inches(6), Inches(0.2),
     "PILOTS, RANKED BY REALIZED VALUE",
     font=FONT_MONO, size=7, color="SLATE", letter_spacing=0.16)
text(s, chart_x + chart_w - Inches(2.5),
     chart_y + chart_h + Inches(0.15), Inches(2.5), Inches(0.2),
     "N ≈ 40 PER BANK", font=FONT_MONO, size=7, color="SLATE",
     align="right", letter_spacing=0.16)

# Figure source line at the bottom inside the box
hline(s, fig_x + Inches(0.3), fig_top + fig_h - Inches(0.4),
      fig_x + fig_w - Inches(0.3), color="RULE", weight=0.5)
text(s, fig_x + Inches(0.3), fig_top + fig_h - Inches(0.32),
     fig_w - Inches(0.6), Inches(0.3),
     "SOURCE  ·  BCG BUILD FOR THE FUTURE 2024  ·  ILLUSTRATIVE SHAPE, IC ANALYSIS",
     font=FONT_MONO, size=7, color="STONE", letter_spacing=0.14)

footer(s, 5)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — The thesis (moneyshot, hero quote)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "The thesis", "Reframe")
chapter_eyebrow(s, "II", "Reframe · Our point of view")

# The big quote — Fraunces, large, italic BLUE on the key phrase
text(s, MARGIN, BODY_TOP + Inches(0.9), BODY_W, Inches(2.8),
     [[("Banks have been adding agents to processes.", {"color": "DEEP"})],
      [("The next decade is about ", {"color": "DEEP"}),
       ("redesigning", {"italic": True, "color": "BLUE"}),
       ("", {})],
      [("processes around agents.", {"italic": True, "color": "BLUE"})]],
     font=FONT_DISPLAY, size=42, color="DEEP", line_spacing=1.15)

# BLUE accent rule — clearly below the three-line quote
rect(s, MARGIN, BODY_TOP + Inches(4.2), Inches(1.4), Inches(0.035), "BLUE")

# Supporting paragraph in Newsreader
text(s, MARGIN, BODY_TOP + Inches(4.55), BODY_W - Inches(2), Inches(1.8),
     "Pilots plateau because the agent inherits the friction of the "
     "workflow it joined — the handoffs, the data silos, the compliance "
     "choke-points, the human approvals built for a slower world. The "
     "teams who break out aren't the ones with better models. They are "
     "the ones who redrew the process around what an agent can actually "
     "do, and rebuilt the operating model — control, risk, talent — to match.",
     font=FONT_BODY, size=14, color="PAVEMENT", line_spacing=1.6)

footer(s, 6)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — Six pillars (2x3 grid, hairlines)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, 'What "agent-native" means in operating-model terms',
           "Reframe  ·  The pillars")
chapter_eyebrow(s, "II", "Pillars")

# Slide subtitle
text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.5),
     "Six pillars define the move from agent-augmented to agent-native.",
     font=FONT_DISPLAY, size=20, italic=True, color="DEEP",
     line_spacing=1.1)

pillars = [
    ("i.",   "Control plane",
     "Orchestration, routing, retries, escalation. The agent's wiring "
     "is a first-class system."),
    ("ii.",  "Agent roles",
     "Each agent has one clear job — research, decide, act, or review. "
     "Roles compose; sprawl is avoided."),
    ("iii.", "Human-in-the-loop",
     "Humans handle exceptions and decisions of consequence. They sit "
     "at the seam, not at the start."),
    ("iv.",  "Observability",
     "Every step is logged, replayable, and scored. What you cannot "
     "measure, you cannot improve."),
    ("v.",   "Governance",
     "Model risk, audit trail, regulator-readable explanations are "
     "designed in, not bolted on."),
    ("vi.",  "Talent",
     "New operating roles: agent designer, agent-ops, agent product "
     "manager. The org chart shifts."),
]

grid_top = BODY_TOP + Inches(1.0)
grid_h = BODY_H - Inches(1.2)
cell_w = (BODY_W - Inches(0.6)) / 3
# Generous gap between rows so the second-row BLUE top rule sits clearly
# below the first row's body text — eliminates the visual "strikethrough"
# the QA pass flagged.
cell_h = (grid_h - Inches(0.7)) / 2
for idx, (roman, head, body) in enumerate(pillars):
    col = idx % 3
    row = idx // 3
    x = MARGIN + col * (cell_w + Inches(0.3))
    y = grid_top + row * (cell_h + Inches(0.7))
    # Thin top BLUE rule
    rect(s, x, y, Inches(0.35), Inches(0.025), "BLUE")
    # Roman numeral
    text(s, x, y + Inches(0.18), Inches(1.5), Inches(0.5), roman,
         font=FONT_DISPLAY, size=22, italic=True, color="BLUE",
         line_spacing=1.0)
    # Head in Fraunces
    text(s, x, y + Inches(0.85), cell_w - Inches(0.2), Inches(0.5), head,
         font=FONT_DISPLAY, size=20, color="DEEP", line_spacing=1.1)
    # Body in Newsreader
    text(s, x, y + Inches(1.5), cell_w - Inches(0.2), Inches(1.4), body,
         font=FONT_BODY, size=12, color="SLATE", line_spacing=1.55)

footer(s, 7)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — The bank, redrawn (3 lanes — NO DEEP_BLUE on headers)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "The bank, redrawn: where the unit of work changes",
           "Reframe  ·  Hot-spots")
chapter_eyebrow(s, "II", "Hot-spots")

text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.6),
     "A finite set of processes are candidates for agent-native "
     "redesign. The blue marks show where the work itself changes shape.",
     font=FONT_BODY, size=13, italic=True, color="SLATE", line_spacing=1.5)

lanes = [
    ("FRONT OFFICE", "Customer, advisor, channel-facing", [
        "Relationship banker copilot",
        "Wealth advisor co-pilot",
        "Digital channel agentic flows",
        "Contact-centre resolution agent",
    ]),
    ("MIDDLE OFFICE", "Risk, ops, control", [
        "KYC & onboarding orchestration",
        "Credit underwriting & memo",
        "Dispute handling & servicing",
        "Fraud investigation co-pilot",
    ]),
    ("BACK OFFICE", "Finance, ops, regulatory", [
        "Reconciliation & break-resolution",
        "Regulatory reporting drafting",
        "Treasury & liquidity ops",
        "Audit & compliance evidence",
    ]),
]
lane_top = BODY_TOP + Inches(1.0)
lane_h = BODY_H - Inches(1.3)
lane_w = (BODY_W - Inches(0.6)) / 3

for i, (head, sub, items) in enumerate(lanes):
    x = MARGIN + i * (lane_w + Inches(0.3))
    # No DEEP_BLUE band — mono header in BLUE, then BLUE seam rule, then content
    text(s, x, lane_top, lane_w, Inches(0.3), head,
         font=FONT_MONO, size=10, bold=True, color="BLUE",
         letter_spacing=0.2)
    # BLUE seam rule below header
    rect(s, x, lane_top + Inches(0.32), lane_w, Inches(0.025), "BLUE")
    # Sub-label in Fraunces italic
    text(s, x, lane_top + Inches(0.45), lane_w, Inches(0.4), sub,
         font=FONT_DISPLAY, size=14, italic=True, color="DEEP",
         line_spacing=1.2)
    # Items with BLUE dot
    iy = lane_top + Inches(1.1)
    for it in items:
        oval(s, x + Inches(0.05), iy + Inches(0.1),
             Inches(0.12), Inches(0.12), "BLUE")
        text(s, x + Inches(0.3), iy, lane_w - Inches(0.4), Inches(0.5),
             it, font=FONT_BODY, size=13, color="PAVEMENT",
             anchor="middle", line_spacing=1.4)
        iy += Inches(0.65)

footer(s, 8)


# ═══════════════════════════════════════════════════════════════════
# SLIDES 8/9/10 — The three patterns (before/after with SEAM motif)
# ═══════════════════════════════════════════════════════════════════
def pattern_slide(n, eyebrow, title, sub_emph, sub_rest, before_lines,
                  after_lines, principles):
    s = prs.slides.add_slide(BLANK)
    title_band(s, title, eyebrow)
    chapter_eyebrow(s, "III", "Three patterns")

    # Sub-headline — Fraunces with italic BLUE on the key phrase
    runs = []
    if sub_emph:
        runs.append((sub_emph, {"italic": True, "color": "BLUE"}))
        if sub_rest:
            runs.append((sub_rest, {"color": "DEEP"}))
    else:
        runs.append((sub_rest, {"color": "DEEP"}))
    text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.6), runs,
         font=FONT_DISPLAY, size=22, color="DEEP", line_spacing=1.1)

    # Two-column before/after with vertical SEAM rule
    col_top = BODY_TOP + Inches(1.0)
    col_h = Inches(3.6)
    col_w = (BODY_W - Inches(0.5)) / 2

    # Before
    bx = MARGIN
    rect(s, bx, col_top, col_w, col_h, "TINT_STONE")
    text(s, bx + Inches(0.3), col_top + Inches(0.25), col_w, Inches(0.3),
         "TODAY", font=FONT_MONO, size=9, bold=True, color="SLATE",
         letter_spacing=0.2)
    rect(s, bx + Inches(0.3), col_top + Inches(0.55),
         Inches(0.4), Inches(0.02), "SLATE")
    text(s, bx + Inches(0.3), col_top + Inches(0.7),
         col_w - Inches(0.6), Inches(0.5), "The legacy process",
         font=FONT_DISPLAY, size=20, italic=True, color="DEEP")
    by = col_top + Inches(1.4)
    for line in before_lines:
        text(s, bx + Inches(0.3), by, Inches(0.15), Inches(0.3),
             "—", font=FONT_BODY, size=12, color="STONE")
        text(s, bx + Inches(0.5), by, col_w - Inches(0.8), Inches(0.5),
             line, font=FONT_BODY, size=12, color="PAVEMENT",
             line_spacing=1.5)
        by += Inches(0.42)

    # SEAM vertical rule between panes
    seam_x = MARGIN + col_w + Inches(0.25)
    rect(s, seam_x - Inches(0.01), col_top, Inches(0.025), col_h, "BLUE")
    # SEAM label
    text(s, seam_x - Inches(0.4), col_top + col_h/2 - Inches(0.15),
         Inches(0.8), Inches(0.3), "SEAM",
         font=FONT_MONO, size=7, bold=True, color="BLUE",
         align="center", anchor="middle", letter_spacing=0.24)

    # After
    ax = MARGIN + col_w + Inches(0.5)
    rect(s, ax, col_top, col_w, col_h, "TINT_BLUE")
    text(s, ax + Inches(0.3), col_top + Inches(0.25), col_w, Inches(0.3),
         "REDESIGNED", font=FONT_MONO, size=9, bold=True, color="BLUE",
         letter_spacing=0.2)
    rect(s, ax + Inches(0.3), col_top + Inches(0.55),
         Inches(0.4), Inches(0.02), "BLUE")
    text(s, ax + Inches(0.3), col_top + Inches(0.7),
         col_w - Inches(0.6), Inches(0.5), "The agent-native version",
         font=FONT_DISPLAY, size=20, italic=True, color="DEEP")
    ay = col_top + Inches(1.4)
    for line in after_lines:
        text(s, ax + Inches(0.3), ay, Inches(0.15), Inches(0.3),
             "+", font=FONT_MONO, size=12, bold=True, color="BLUE")
        text(s, ax + Inches(0.5), ay, col_w - Inches(0.8), Inches(0.5),
             line, font=FONT_BODY, size=12, color="PAVEMENT",
             line_spacing=1.5)
        ay += Inches(0.42)

    # Three redesign principles below the panes
    p_top = col_top + col_h + Inches(0.35)
    text(s, MARGIN, p_top, BODY_W, Inches(0.3),
         "REDESIGN PRINCIPLES",
         font=FONT_MONO, size=9, bold=True, color="BLUE",
         letter_spacing=0.22)
    p_top2 = p_top + Inches(0.35)
    pw = (BODY_W - Inches(0.6)) / 3
    for i, p in enumerate(principles):
        px = MARGIN + i * (pw + Inches(0.3))
        # Thin BLUE left rule
        rect(s, px, p_top2, Inches(0.03), Inches(0.7), "BLUE")
        text(s, px + Inches(0.18), p_top2, pw - Inches(0.2),
             Inches(0.85), p, font=FONT_BODY, size=11, italic=True,
             color="DEEP", line_spacing=1.45)

    footer(s, n)


pattern_slide(
    9, "Proof  ·  Pattern 01", "Onboarding & KYC",
    sub_emph=None,
    sub_rest="14 systems, 11 handoffs, 9 days — and the analyst spent most of it moving documents.",
    before_lines=[
        "~14 systems, 11 handoffs",
        "9-day median TAT for commercial",
        "8–12% drop-off in customer journey",
        "Analyst-led, document-by-document review",
    ],
    after_lines=[
        "Single orchestrated agent flow",
        "Documents read, profiles drafted, risk tiered",
        "Humans approve only mid/high-risk cases",
        "Audit trail generated continuously",
    ],
    principles=[
        "The customer enters the system once. The system does the joining.",
        "Documents go to an agent first, an analyst second.",
        "Control sits in the orchestration layer, not the screens.",
    ],
)

pattern_slide(
    10, "Proof  ·  Pattern 02", "Credit underwriting & decisioning",
    sub_emph="The memo is generated, not authored.",
    sub_rest=" The reviewer sees the working.",
    before_lines=[
        "Analyst pulls statements, spreads, market data",
        "Manual memo drafting over weeks",
        "Limited scenario coverage (2–3 cases)",
        "Reviewer sees the answer, rarely the working",
    ],
    after_lines=[
        "Research agent ingests filings + market data",
        "Decision-support agent runs spread and scenarios",
        "Human reviewer at the seam, not the start",
        "Memo generated from the agent's working",
    ],
    principles=[
        "The memo is generated, not authored.",
        "Scenarios are cheap; run more of them.",
        "Reviewers see the working, not just the answer.",
    ],
)

pattern_slide(
    11, "Proof  ·  Pattern 03", "Dispute handling & servicing",
    sub_emph="The agent owns the case.",
    sub_rest=" The human owns the exception.",
    before_lines=[
        "25-day median resolution cycle",
        "Customer re-tells story to 2–3 teams",
        "Ledger entries lag the case state",
        "Communications are ticket-driven",
    ],
    after_lines=[
        "One agent owns the case end-to-end",
        "Ledger is the source of truth; screens are the lens",
        "Humans handle exceptions only",
        "Communications generated proactively",
    ],
    principles=[
        "The agent owns the case; the human owns the exception.",
        "The ledger is the truth; the screen is the lens.",
        "Communications are continuous, not on request.",
    ],
)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — What the three patterns share (5 principles)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "What every agent-native redesign has in common",
           "Proof  ·  Synthesis")
chapter_eyebrow(s, "III", "Synthesis")

text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.6),
     "Five recurring principles, regardless of which process you start with.",
     font=FONT_DISPLAY, size=20, italic=True, color="DEEP", line_spacing=1.1)

principles = [
    ("I",   "Events, not handoffs.",
     "The process becomes an event stream. State is shared, not passed between teams."),
    ("II",  "Verify, don't trust.",
     "Every agent output is scored against an eval before it advances. No eval, no deployment."),
    ("III", "Humans at the seam.",
     "Humans handle exceptions and decisions of consequence — not data movement or formatting."),
    ("IV",  "Ledger first, screen second.",
     "The system of record drives the experience. The screen is a view, not the truth."),
    ("V",   "One agent, many tools.",
     "Agents specialize by role, not by screen. Capability composes; sprawl is avoided."),
]
list_top = BODY_TOP + Inches(0.9)
row_h = Inches(0.95)
for i, (roman, head, body) in enumerate(principles):
    y = list_top + i * row_h
    # Hairline rule
    hline(s, MARGIN, y, MARGIN + BODY_W, color="RULE", weight=0.5)
    # Roman in Fraunces italic BLUE
    text(s, MARGIN, y + Inches(0.15), Inches(1.0), Inches(0.8), roman,
         font=FONT_DISPLAY, size=44, italic=True, color="BLUE",
         line_spacing=1.0)
    # Head + body
    text(s, MARGIN + Inches(1.2), y + Inches(0.18),
         Inches(4.6), Inches(0.6), head,
         font=FONT_DISPLAY, size=20, color="DEEP", line_spacing=1.15)
    text(s, MARGIN + Inches(5.9), y + Inches(0.2),
         BODY_W - Inches(5.9), Inches(0.7), body,
         font=FONT_BODY, size=13, color="SLATE", line_spacing=1.5)
# Closing rule
hline(s, MARGIN, list_top + 5*row_h, MARGIN + BODY_W,
      color="RULE", weight=0.5)

footer(s, 12)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — The agentic maturity model (horizontal stepped axis)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "The agentic maturity model",
           "Maturity & path  ·  Four stages")
chapter_eyebrow(s, "IV", "Maturity")

text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.6),
     "Maturity is not a model question. It is an operating-model question.",
     font=FONT_DISPLAY, size=22, italic=True, color="DEEP", line_spacing=1.1)

# Horizontal axis — BLUE rule across. Pulled in from both edges so that
# each stage's column of text (3in wide centred on its dot) stays inside
# the slide's safe area without clipping.
axis_y = BODY_TOP + Inches(1.9)
axis_x1 = MARGIN + Inches(1.5)
axis_x2 = SLIDE_W - MARGIN - Inches(1.5)
hline(s, axis_x1, axis_y, axis_x2, color="BLUE", weight=1.5)

stages = [
    ("STAGE I",   "Today",   "Pilots",
     "One-off use cases. Manual scaling. Value rarely attributed back to a P&L.", False),
    ("STAGE II",  "Today",   "Patterned use cases",
     "Templates emerge. Platform investments begin. Evals are introduced.", False),
    ("STAGE III", "The prize", "Agent-native processes",
     "Workflows redesigned end-to-end. KPIs rewritten. Operating model adjusted.", True),
    ("STAGE IV",  "The future", "Agent-native bank",
     "Roles, P&L, and controls are built around agents. The org chart reflects it.", True),
]
n = len(stages)
seg = (axis_x2 - axis_x1) / (n - 1)
for i, (stage_no, tag, name, body, future) in enumerate(stages):
    cx = axis_x1 + i * seg
    # Dot
    if future:
        oval(s, cx - Inches(0.1), axis_y - Inches(0.1),
             Inches(0.2), Inches(0.2), "WHITE", line="BLUE", line_w=1.5)
    else:
        oval(s, cx - Inches(0.1), axis_y - Inches(0.1),
             Inches(0.2), Inches(0.2), "BLUE")
    # Stage label above
    text(s, cx - Inches(1.5), axis_y - Inches(1.0),
         Inches(3.0), Inches(0.3), stage_no,
         font=FONT_MONO, size=9, bold=True, color="BLUE",
         align="center", letter_spacing=0.2)
    text(s, cx - Inches(1.5), axis_y - Inches(0.7),
         Inches(3.0), Inches(0.3), f"· {tag.upper()} ·",
         font=FONT_MONO, size=7,
         color="STONE" if not future else "BLUE",
         align="center", letter_spacing=0.18)
    # Stage name below
    text(s, cx - Inches(1.5), axis_y + Inches(0.3),
         Inches(3.0), Inches(0.5), name,
         font=FONT_DISPLAY, size=18, color="DEEP",
         align="center", line_spacing=1.15)
    # Stage description
    text(s, cx - Inches(1.5), axis_y + Inches(1.0),
         Inches(3.0), Inches(1.8), body,
         font=FONT_BODY, size=11, color="SLATE",
         align="center", line_spacing=1.5)

footer(s, 13)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — Where banks are (distribution chart)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "Where most banks are — and where leaders are heading",
           "Maturity & path  ·  Today's distribution")
chapter_eyebrow(s, "IV", "Today")

text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.6),
     "Most large banks remain between stages one and two. The gap to "
     "stage three is where the next decade of advantage is decided.",
     font=FONT_DISPLAY, size=20, italic=True, color="DEEP", line_spacing=1.2)

# Distribution chart — 4 bars
chart_top = BODY_TOP + Inches(1.4)
chart_h = Inches(2.6)
chart_x = MARGIN + Inches(1.0)
chart_w = BODY_W - Inches(2.0)
seg_w = chart_w / 4
labels = ["Pilots", "Patterned\nuse cases",
          "Agent-native\nprocesses", "Agent-native\nbank"]
romans = ["I", "II", "III", "IV"]
populations = [0.40, 0.42, 0.16, 0.02]

# Baseline
hline(s, chart_x, chart_top + chart_h, chart_x + chart_w,
      color="PAVEMENT", weight=0.75)
max_bar_h = chart_h - Inches(0.3)
# All percentage labels sit at a SHARED Y above the chart top, so they read
# as a consistent label row regardless of bar height. (Otherwise the tiny
# 2% bar's label floats next to its baseline while the 40%/42% labels sit
# at the top — visually inconsistent.)
label_y = chart_top - Inches(0.05)
for i, (lab, pop) in enumerate(zip(labels, populations)):
    bx = chart_x + i * seg_w + Inches(0.5)
    bw = seg_w - Inches(1.0)
    bh = Emu(int(max_bar_h * (pop / 0.45)))
    by = chart_top + chart_h - bh
    color = "BLUE" if i >= 2 else "STONE"
    rect(s, bx, by, bw, bh, color)
    # Percentage label — large Fraunces, anchored at a shared label row
    text(s, bx - Inches(0.3), label_y, bw + Inches(0.6),
         Inches(0.55), f"{int(pop*100)}%",
         font=FONT_DISPLAY, size=32, italic=False,
         color="DEEP" if i < 2 else "BLUE", align="center",
         line_spacing=1.0)
    # Roman numeral STAGE marker below — single line, mono
    text(s, bx - Inches(0.4), chart_top + chart_h + Inches(0.15),
         bw + Inches(0.8), Inches(0.25),
         f"STAGE  {romans[i]}",
         font=FONT_MONO, size=9, bold=True,
         color="BLUE" if i >= 2 else "STONE",
         align="center", letter_spacing=0.22)
    # Stage name (multi-line) below the Roman marker
    text(s, bx - Inches(0.4), chart_top + chart_h + Inches(0.5),
         bw + Inches(0.8), Inches(0.7), lab,
         font=FONT_BODY, size=12, italic=True, color="PAVEMENT",
         align="center", line_spacing=1.3)

# Insight callout — positioned with clear space below labels
text(s, MARGIN, BODY_BOT - Inches(0.95), BODY_W, Inches(0.5),
     "The gap between stages II and III is where the next decade is decided.",
     font=FONT_DISPLAY, size=16, italic=True, color="BLUE",
     line_spacing=1.3)
# Source line — pinned to the very bottom inside the body area
text(s, MARGIN, BODY_BOT - Inches(0.35), BODY_W, Inches(0.25),
     "SOURCE  ·  IC ANALYSIS OF PUBLIC DISCLOSURES, MCKINSEY STATE OF AI 2024, BCG 2024  ·  ILLUSTRATIVE",
     font=FONT_MONO, size=7, color="STONE", letter_spacing=0.14)

footer(s, 14)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — The path to value (3 horizons as horizontal rows)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "The path to value: what to do across 24 months",
           "Maturity & path  ·  24-month plan")
chapter_eyebrow(s, "IV", "The path")

text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.5),
     "Three horizons. Leaders cross from pilots to a live agent-native "
     "process inside this window.",
     font=FONT_DISPLAY, size=18, italic=True, color="DEEP", line_spacing=1.2)

horizons = [
    ("0 – 3 MONTHS", "Build the case", [
        "Pick 1–2 candidate processes for redesign",
        "Stand up evals + governance",
        "Quantify the value at stake; secure executive sponsor",
        "Recruit / re-skill the core agent team",
    ]),
    ("3 – 12 MONTHS", "Ship one process", [
        "Redesign and deliver one process end-to-end",
        "Capture and attribute savings",
        "Codify reusable patterns",
        "Begin shifting KPIs around the new process",
    ]),
    ("12 – 24 MONTHS", "Scale the model", [
        "Roll patterns across the function",
        "Restructure ops around agent-native flows",
        "Rewrite productivity and risk KPIs",
        "Stand up agent-ops as a permanent capability",
    ]),
]
top = BODY_TOP + Inches(1.0)
row_h = (BODY_H - Inches(1.3)) / 3
for i, (when, verb, items) in enumerate(horizons):
    y = top + i * row_h
    # Hairline rule above
    hline(s, MARGIN, y, MARGIN + BODY_W, color="RULE", weight=0.5)
    # Left column: mono horizon + italic Fraunces verb
    text(s, MARGIN, y + Inches(0.2), Inches(2.2), Inches(0.3),
         when, font=FONT_MONO, size=9, bold=True, color="BLUE",
         letter_spacing=0.18)
    text(s, MARGIN, y + Inches(0.55), Inches(2.4), Inches(0.6),
         verb, font=FONT_DISPLAY, size=22, italic=True, color="DEEP",
         line_spacing=1.15)
    # Right column: items in 2-column layout
    ix1 = MARGIN + Inches(3.0)
    icw = (BODY_W - Inches(3.0)) / 2
    for j, it in enumerate(items):
        col = j % 2
        row = j // 2
        ix = ix1 + col * icw
        iy = y + Inches(0.15) + row * Inches(0.55)
        text(s, ix, iy, Inches(0.15), Inches(0.4),
             "—", font=FONT_BODY, size=12, color="BLUE")
        text(s, ix + Inches(0.2), iy, icw - Inches(0.3), Inches(0.5),
             it, font=FONT_BODY, size=12, color="PAVEMENT",
             line_spacing=1.45)
# Closing rule
hline(s, MARGIN, top + 3*row_h, MARGIN + BODY_W, color="RULE", weight=0.5)

footer(s, 15)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 15 — Five no-regret moves (Roman numerals, editorial)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "Five no-regret moves to start now",
           "Maturity & path  ·  Moves")
chapter_eyebrow(s, "IV", "Moves")

text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.5),
     "Infrastructural, not speculative. They pay back in every plausible "
     "future state of the technology.",
     font=FONT_DISPLAY, size=18, italic=True, color="DEEP", line_spacing=1.2)

moves = [
    ("i",   "Establish an agent control plane.",
     "Pick one orchestration stack; standardize. Multiple stacks fragment "
     "evals, controls, and talent."),
    ("ii",  "Build the eval discipline.",
     "Every agentic capability ships with a measurable scorecard. No "
     "eval, no promotion."),
    ("iii", "Pick a redesign, not a copilot.",
     "Choose a process worth redrawing — not just augmenting. Copilots "
     "have known ceilings; redesigns don't."),
    ("iv",  "Get the governance ready early.",
     "Model risk, audit trail, regulator-readable explanations from day "
     "one. First-movers shape the precedent."),
    ("v",   "Reshape the talent contract.",
     "Hire and train agent designers, agent-ops, and agent product "
     "managers. The role gap closes in 6–12 months."),
]
top = BODY_TOP + Inches(1.0)
row_h = (BODY_H - Inches(1.2)) / 5
for i, (roman, head, body) in enumerate(moves):
    y = top + i * row_h
    hline(s, MARGIN, y, MARGIN + BODY_W, color="RULE", weight=0.5)
    text(s, MARGIN, y + Inches(0.12), Inches(1.2), Inches(0.7), roman,
         font=FONT_DISPLAY, size=28, italic=True, color="BLUE",
         line_spacing=1.0)
    text(s, MARGIN + Inches(1.2), y + Inches(0.15),
         BODY_W - Inches(1.2), Inches(0.45), head,
         font=FONT_DISPLAY, size=18, color="DEEP", line_spacing=1.15)
    text(s, MARGIN + Inches(1.2), y + Inches(0.55),
         BODY_W - Inches(1.2), Inches(0.45), body,
         font=FONT_BODY, size=12, color="SLATE", line_spacing=1.5)
hline(s, MARGIN, top + 5*row_h, MARGIN + BODY_W, color="RULE", weight=0.5)

footer(s, 16)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 16 — IC capability map [PLACEHOLDER]
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "How Infosys Consulting helps",
           "IC  ·  Capability map")
chapter_eyebrow(s, "V", "How IC helps")
placeholder_pill(s, "Replace with IC official copy")

text(s, MARGIN, BODY_TOP + Inches(0.6), BODY_W, Inches(0.5),
     "Four capabilities, across the agentic lifecycle.",
     font=FONT_DISPLAY, size=20, italic=True, color="DEEP", line_spacing=1.1)

caps = [
    ("STRATEGY", "Shape the bet",
     ["Agentic AI strategy & value-pool mapping",
      "Operating-model design",
      "Investment case & ROI framing"]),
    ("DESIGN", "Redraw the process",
     ["Process redesign for agent-native flows",
      "Agent role & control-plane design",
      "Governance & risk framework"]),
    ("BUILD", "Make it real",
     ["Agent engineering & orchestration",
      "Integration with core banking",
      "Evals, observability, guardrails"]),
    ("RUN", "Keep it improving",
     ["Agent-ops capability stand-up",
      "Model risk & continuous monitoring",
      "Capability rollout & change"]),
]
grid_top = BODY_TOP + Inches(1.4)
grid_h = BODY_H - Inches(1.6)
cell_w = (BODY_W - Inches(0.6)) / 4
for i, (eyebrow, head, items) in enumerate(caps):
    x = MARGIN + i * (cell_w + Inches(0.2))
    # Mono header in BLUE (NO DEEP_BLUE band)
    text(s, x, grid_top, cell_w, Inches(0.3), eyebrow,
         font=FONT_MONO, size=10, bold=True, color="BLUE",
         letter_spacing=0.2)
    # BLUE seam under the header
    rect(s, x, grid_top + Inches(0.32), cell_w, Inches(0.025), "BLUE")
    # Fraunces head
    text(s, x, grid_top + Inches(0.45), cell_w, Inches(0.5), head,
         font=FONT_DISPLAY, size=18, italic=True, color="DEEP",
         line_spacing=1.15)
    # Items
    iy = grid_top + Inches(1.2)
    for it in items:
        text(s, x, iy, Inches(0.15), Inches(0.3), "—",
             font=FONT_BODY, size=11, color="BLUE")
        text(s, x + Inches(0.2), iy, cell_w - Inches(0.3), Inches(0.7),
             it, font=FONT_BODY, size=11, color="PAVEMENT",
             line_spacing=1.5)
        iy += Inches(0.7)

footer(s, 17)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 17 — Where IC plays in the stack [PLACEHOLDER]
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "Where IC plays across the agentic stack", "IC  ·  Stack")
chapter_eyebrow(s, "V", "Stack")
placeholder_pill(s, "Confirm IC positioning")

text(s, MARGIN, BODY_TOP + Inches(0.6), BODY_W, Inches(0.5),
     "The agent stack — and the layers where IC delivers.",
     font=FONT_DISPLAY, size=20, italic=True, color="DEEP", line_spacing=1.1)

layers = [
    ("Change management & adoption", True,
     "Operating-model change, role redesign, learning"),
    ("Governance, model risk & audit", True,
     "Risk framework, audit trail, regulator readiness"),
    ("Agent operations", True,
     "Run, monitor, improve, retire — the lifecycle"),
    ("Integration with core banking", True,
     "Connect agents to systems of record"),
    ("Agent design & orchestration", True,
     "Roles, control plane, evals, guardrails"),
    ("Tooling & infrastructure", False,
     "Frameworks, vector stores, observability"),
    ("Foundation models", False,
     "Models from providers (selection advisory only)"),
]
stack_top = BODY_TOP + Inches(1.4)
stack_h = BODY_H - Inches(1.6)
layer_h = stack_h / len(layers)
for i, (name, ic, desc) in enumerate(layers):
    y = stack_top + i * layer_h
    # Hairline rule top
    hline(s, MARGIN, y, MARGIN + BODY_W, color="RULE", weight=0.5)
    # Left rule strip — BLUE for IC plays, STONE for advisory only.
    # Thin vertical strip, NOT a full-width band.
    rule_color = "BLUE" if ic else "STONE"
    rect(s, MARGIN, y, Inches(0.05), layer_h - Inches(0.04), rule_color)
    # Status tag inside a small pill, left-positioned
    tag = "IC PLAYS HERE" if ic else "ADVISORY ONLY"
    tag_color = "BLUE" if ic else "STONE"
    tag_fill = "TINT_BLUE" if ic else "TINT_STONE"
    rect(s, MARGIN + Inches(0.25), y + (layer_h - Inches(0.34))/2,
         Inches(1.7), Inches(0.28), tag_fill)
    text(s, MARGIN + Inches(0.25), y + (layer_h - Inches(0.34))/2,
         Inches(1.7), Inches(0.28), tag,
         font=FONT_MONO, size=7, bold=True, color=tag_color,
         align="center", anchor="middle", letter_spacing=0.18)
    # Layer name — Fraunces italic
    text(s, MARGIN + Inches(2.3), y, Inches(5.2),
         layer_h - Inches(0.04), name,
         font=FONT_DISPLAY, size=16, italic=True, color="DEEP",
         anchor="middle")
    # Description
    text(s, MARGIN + Inches(7.8), y, BODY_W - Inches(8.0),
         layer_h - Inches(0.04), desc,
         font=FONT_BODY, size=11, color="SLATE",
         anchor="middle", line_spacing=1.4)
# Closing hairline rule
hline(s, MARGIN, stack_top + len(layers)*layer_h,
      MARGIN + BODY_W, color="RULE", weight=0.5)

footer(s, 18)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 18 — Qualifying questions (2x2 grid, toolbox style)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "Qualifying questions for client conversations",
           "Sales tool  ·  In the room")
chapter_eyebrow(s, "V", "Tool 01")

text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.4),
     "Ten questions, four dimensions. The answers reveal where the client sits.",
     font=FONT_DISPLAY, size=17, italic=True, color="DEEP", line_spacing=1.2)

groups = [
    ("AMBITION", [
        "What process would you redesign if you could start over today?",
        "If agents become 50% cheaper next year, what changes in your operating model?",
    ]),
    ("CURRENT STATE", [
        "How many of your pilots have measurable value attribution today?",
        "Which use cases are in production with daily users?",
        "Which use cases were quietly retired in the last 12 months?",
    ]),
    ("RISK APPETITE", [
        "What level of agent autonomy is your regulator comfortable with today?",
        "Which decisions must always have a human reviewer — and why?",
    ]),
    ("OPERATING-MODEL READINESS", [
        "Who owns 'agent ops' in your organization?",
        "How will you measure productivity in an agent-native team?",
        "Where will the savings show up — and who is accountable for them?",
    ]),
]
# Move groups up; widen the available grid so 3-question groups fit cleanly.
grid_top = BODY_TOP + Inches(0.75)
grid_h = BODY_H - Inches(0.9)
gw = (BODY_W - Inches(0.6)) / 2
gh = (grid_h - Inches(0.4)) / 2
for idx, (head, qs) in enumerate(groups):
    col = idx % 2
    row = idx // 2
    x = MARGIN + col * (gw + Inches(0.6))
    y = grid_top + row * (gh + Inches(0.4))
    # Mono header in BLUE
    text(s, x, y, gw, Inches(0.3), head,
         font=FONT_MONO, size=9, bold=True, color="BLUE",
         letter_spacing=0.22)
    # BLUE seam
    rect(s, x, y + Inches(0.32), Inches(0.4), Inches(0.025), "BLUE")
    # Questions — size 11 keeps long questions on a single line at this width,
    # and 0.65in pitch leaves clearance if any wraps.
    qy = y + Inches(0.5)
    row_pitch = Inches(0.65)
    for q in qs:
        text(s, x, qy, Inches(0.35), Inches(0.4), "Q.",
             font=FONT_MONO, size=10, bold=True, color="BLUE")
        text(s, x + Inches(0.35), qy, gw - Inches(0.45), row_pitch - Inches(0.05),
             q, font=FONT_BODY, size=11, italic=True, color="PAVEMENT",
             line_spacing=1.4)
        qy += row_pitch

footer(s, 19)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 19 — Common objections
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "Common objections — and how to answer them",
           "Sales tool  ·  Objections")
chapter_eyebrow(s, "V", "Tool 02")

text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.5),
     "Five things you will hear in the first 30 minutes — and how to answer each.",
     font=FONT_DISPLAY, size=18, italic=True, color="DEEP", line_spacing=1.2)

rows = [
    ("Our regulator will never allow autonomous agents.",
     "Start with agent-assisted, not autonomous. Design the audit trail before the agent. Regulators react to surprises, not to readable explanations."),
    ("We need to wait for the tech to stabilize.",
     "The agent layer changes; the redesign value is durable. Build patterns, not POCs."),
    ("We've tried — pilots don't scale.",
     "Pilots fail because the process is unchanged. Redesign the process, then scale. The plateau is a design problem, not a model problem."),
    ("Our data isn't ready.",
     "Agents tolerate messier data than traditional ML. Start with read-mostly use cases. Data investment runs in parallel."),
    ("We don't have the talent.",
     "Talent reshapes faster than tech. The role gap is solvable in 6–12 months with a mix of hires and re-skilling."),
]
top = BODY_TOP + Inches(1.0)
row_h = (BODY_H - Inches(1.2)) / len(rows)
for i, (said, resp) in enumerate(rows):
    y = top + i * row_h
    hline(s, MARGIN, y, MARGIN + BODY_W, color="RULE", weight=0.5)
    # Said: italic Fraunces with BLUE quotes
    text(s, MARGIN, y + Inches(0.15),
         BODY_W * 0.42, row_h - Inches(0.3),
         [[("“", {"color": "BLUE", "italic": False}),
           (said, {"italic": True, "color": "DEEP"}),
           ("”", {"color": "BLUE", "italic": False})]],
         font=FONT_DISPLAY, size=16, color="DEEP",
         anchor="middle", line_spacing=1.3)
    # Arrow indicator
    text(s, MARGIN + BODY_W * 0.44, y + Inches(0.2),
         Inches(0.3), Inches(0.5),
         "↳", font=FONT_MONO, size=14, bold=True, color="BLUE",
         anchor="middle")
    # Response in Newsreader
    text(s, MARGIN + BODY_W * 0.48, y + Inches(0.15),
         BODY_W * 0.52, row_h - Inches(0.3), resp,
         font=FONT_BODY, size=12, color="PAVEMENT",
         anchor="middle", line_spacing=1.55)
hline(s, MARGIN, top + len(rows)*row_h, MARGIN + BODY_W,
      color="RULE", weight=0.5)

footer(s, 20)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 20 — Illustrative case studies [PLACEHOLDER]
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "Illustrative case studies", "IC  ·  Proof points")
chapter_eyebrow(s, "V", "Cases")
placeholder_pill(s, "Replace with real engagements")

cases = [
    ("Top-10 US bank", "Agent-native KYC",
     [("Situation", "14-system onboarding, 9-day TAT, 8% drop-off"),
      ("Redesign",  "Orchestrated agent flow with risk-tiered checkpoint"),
      ("Outcome",   "60%+ TAT reduction, 40% analyst capacity returned")]),
    ("European universal bank", "Credit underwriting copilot",
     [("Situation", "Multi-week memo cycle, 2–3 scenarios per deal"),
      ("Redesign",  "Research + decision-support agents, reviewer at the seam"),
      ("Outcome",   "35% memo turnaround reduction, 2× scenario coverage")]),
    ("Regional commercial bank", "Dispute handling redesign",
     [("Situation", "25-day cycle, customer re-tells story to multiple teams"),
      ("Redesign",  "Single agent owns case end-to-end, ledger-first"),
      ("Outcome",   "25-day → 7-day median resolution, +18 NPS lift")]),
]
top = BODY_TOP + Inches(0.7)
card_w = (BODY_W - Inches(0.6)) / 3
card_h = BODY_H - Inches(0.9)
for i, (client, project, rows) in enumerate(cases):
    x = MARGIN + i * (card_w + Inches(0.3))
    # Top thin BLUE rule
    rect(s, x, top, Inches(0.5), Inches(0.025), "BLUE")
    # Client (mono eyebrow)
    text(s, x, top + Inches(0.15), card_w, Inches(0.3),
         client.upper(),
         font=FONT_MONO, size=9, bold=True, color="BLUE",
         letter_spacing=0.18)
    # Project (Fraunces italic)
    text(s, x, top + Inches(0.5), card_w, Inches(0.9), project,
         font=FONT_DISPLAY, size=20, italic=True, color="DEEP",
         line_spacing=1.15)
    # Rows — Situation / Redesign / Outcome
    ry = top + Inches(1.7)
    for label, body in rows:
        text(s, x, ry, card_w, Inches(0.25), label.upper(),
             font=FONT_MONO, size=8, bold=True, color="SLATE",
             letter_spacing=0.18)
        text(s, x, ry + Inches(0.28), card_w - Inches(0.2),
             Inches(0.9), body,
             font=FONT_BODY, size=12, color="PAVEMENT", line_spacing=1.5)
        ry += Inches(1.1)

# Single illustrative-tag note pinned to the very bottom of the slide
text(s, MARGIN, BODY_BOT - Inches(0.25), BODY_W, Inches(0.2),
     "ALL THREE CASES ILLUSTRATIVE  ·  REPLACE WITH REDACTED REAL IC ENGAGEMENTS BEFORE EXTERNAL USE",
     font=FONT_MONO, size=7, color="STONE", align="center",
     letter_spacing=0.16)

footer(s, 21)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 21 — The bottom line (close)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "The bottom line", "Close")
chapter_eyebrow(s, "VI", "Close · What to take away")

# Big Fraunces statement with BLUE italic emphasis — pulled up; the eyebrow
# already labels the section.
text(s, MARGIN, BODY_TOP + Inches(0.6), BODY_W, Inches(2.4),
     [[("Stop adding agents.", {"color": "DEEP"})],
      [("Start ", {"color": "DEEP"}),
       ("redesigning", {"italic": True, "color": "BLUE"}),
       (" around them.", {"color": "DEEP"})]],
     font=FONT_DISPLAY, size=46, color="DEEP", line_spacing=1.18)

# BLUE accent rule directly beneath the statement
rect(s, MARGIN, BODY_TOP + Inches(3.15), Inches(1.2), Inches(0.035), "BLUE")

# Three numbered takeaways
takeaways = [
    "01", "The pilot plateau is a design problem, not a model problem.",
    "02", "Value lives in process redesign and operating-model change — not in better assistants.",
    "03", "The next 24 months will separate banks that 'use AI' from banks that are agent-native.",
]
ty = BODY_TOP + Inches(3.6)
for i in range(3):
    text(s, MARGIN, ty + Inches(0.05), Inches(0.7), Inches(0.4),
         takeaways[i*2],
         font=FONT_MONO, size=10, bold=True, color="BLUE",
         letter_spacing=0.18)
    text(s, MARGIN + Inches(0.7), ty, BODY_W - Inches(0.7), Inches(0.55),
         takeaways[i*2 + 1],
         font=FONT_BODY, size=14, italic=False, color="PAVEMENT",
         line_spacing=1.5)
    ty += Inches(0.6)

footer(s, 22)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 22 — Next steps & contact [PLACEHOLDER]
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "Where we can take this", "Close  ·  Next conversation")
chapter_eyebrow(s, "VI", "Next")
placeholder_pill(s, "Confirm IC contact details")

# Two columns
col_w = (BODY_W - Inches(0.6)) / 2
col_top = BODY_TOP + Inches(0.7)

# LEFT — suggested actions
text(s, MARGIN, col_top, col_w, Inches(0.3),
     "SUGGESTED NEXT STEPS",
     font=FONT_MONO, size=9, bold=True, color="BLUE",
     letter_spacing=0.22)
rect(s, MARGIN, col_top + Inches(0.32), Inches(0.4), Inches(0.025), "BLUE")

steps = [
    ("Operating-model workshop",
     "A 2-day session with the executive sponsor to identify the first redesign candidate."),
    ("Maturity assessment",
     "A 3-week diagnostic mapping current state against the four-stage maturity model."),
    ("Value pool diagnostic",
     "A 4–6 week analysis of where agent-native redesign moves the needle most."),
]
sy = col_top + Inches(0.6)
for head, body in steps:
    text(s, MARGIN, sy, col_w, Inches(0.4), head,
         font=FONT_DISPLAY, size=16, italic=True, color="DEEP",
         line_spacing=1.15)
    text(s, MARGIN, sy + Inches(0.4), col_w, Inches(0.6), body,
         font=FONT_BODY, size=12, color="SLATE", line_spacing=1.5)
    sy += Inches(1.15)

# RIGHT — contact card
rx = MARGIN + col_w + Inches(0.6)
text(s, rx, col_top, col_w, Inches(0.3), "TALK TO US",
     font=FONT_MONO, size=9, bold=True, color="BLUE",
     letter_spacing=0.22)
rect(s, rx, col_top + Inches(0.32), Inches(0.4), Inches(0.025), "BLUE")

card_top = col_top + Inches(0.6)
card_h = Inches(3.5)
rect(s, rx, card_top, col_w, card_h, "TINT_BLUE")

# Name in Fraunces
text(s, rx + Inches(0.3), card_top + Inches(0.4),
     col_w - Inches(0.6), Inches(0.5), "[Name]",
     font=FONT_DISPLAY, size=22, color="DEEP")
# Title
text(s, rx + Inches(0.3), card_top + Inches(0.95),
     col_w - Inches(0.6), Inches(0.3),
     "[Title — Banking practice, Infosys Consulting]",
     font=FONT_BODY, size=12, italic=True, color="SLATE")
# Divider
rect(s, rx + Inches(0.3), card_top + Inches(1.4),
     Inches(0.7), Inches(0.025), "BLUE")
# Contact lines
text(s, rx + Inches(0.3), card_top + Inches(1.6),
     col_w - Inches(0.6), Inches(0.4),
     "[email@infosysconsulting.com]",
     font=FONT_MONO, size=11, color="BLUE")
text(s, rx + Inches(0.3), card_top + Inches(2.0),
     col_w - Inches(0.6), Inches(0.4),
     "[+1 ___ ___ ____]",
     font=FONT_MONO, size=11, color="SLATE")
text(s, rx + Inches(0.3), card_top + Inches(2.7),
     col_w - Inches(0.6), Inches(0.6),
     "Or ask your IC account team to set up a working session.",
     font=FONT_BODY, size=11, italic=True, color="SLATE", line_spacing=1.5)

footer(s, 23)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 24 — Appendix separator (full-bleed DEEP_BLUE divider)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
# Full-bleed DEEP_BLUE — the deck's only mid-body separator slide
rect(s, 0, 0, SLIDE_W, SLIDE_H, "DEEP")

# Top masthead in mono — same scaffolding as the cover
text(s, MARGIN, Inches(0.55), Inches(8), Inches(0.3),
     "INFOSYS  CONSULTING  ·  AGENTIC AI IN BANKING",
     font=FONT_MONO, size=10, bold=True, color="LT_BLUE",
     letter_spacing=0.2)
text(s, SLIDE_W - MARGIN - Inches(5), Inches(0.55), Inches(5), Inches(0.3),
     "POINT OF VIEW  ·  APPENDIX",
     font=FONT_MONO, size=9, color="STONE",
     align="right", letter_spacing=0.18)

# Chapter mark (centred-ish, large mono kicker)
text(s, MARGIN, Inches(2.9), BODY_W, Inches(0.35),
     "—  CHAPTER  END",
     font=FONT_MONO, size=11, bold=True, color="LT_BLUE",
     letter_spacing=0.28)

# Huge Fraunces "Appendix · Sources and references"
text(s, MARGIN, Inches(3.4), BODY_W, Inches(2.6),
     [[("Appendix.", {"color": "WHITE"})],
      [("Sources ", {"color": "WHITE"}),
       ("&", {"italic": True, "color": "LT_BLUE"}),
       (" references.", {"italic": True, "color": "WHITE"})]],
     font=FONT_DISPLAY, size=76, color="WHITE", line_spacing=1.06)

# Bottom thin BLUE seam — closes the divider
rect(s, 0, SLIDE_H - Inches(0.05), SLIDE_W, Inches(0.05), "BLUE")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 25 — Sources & references
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
title_band(s, "Sources & references", "Appendix")
chapter_eyebrow(s, "—", "Appendix")

text(s, MARGIN, BODY_TOP + Inches(0.05), BODY_W, Inches(0.5),
     "Supporting research and citations across the deck.",
     font=FONT_DISPLAY, size=18, italic=True, color="DEEP")

sources = [
    ("MARKET SIZING & MATURITY", [
        "McKinsey & Co. The economic potential of generative AI, June 2023.",
        "McKinsey & Co. The state of AI in early 2024 — global survey.",
        "Boston Consulting Group. Build for the Future 2024.",
        "Gartner. GenAI implementation survey 2024 — pilot abandonment.",
    ]),
    ("BANKING-SPECIFIC EVIDENCE", [
        "JPMorgan Chase. 2024 Annual Report and 10-K filings.",
        "Bank of America. Investor day disclosures (Erica interactions).",
        "Citigroup. Public statements on agentic AI initiatives, 2024-2025.",
        "Goldman Sachs. Public disclosures on internal AI deployments.",
    ]),
    ("PROCESS BENCHMARKS", [
        "Celent. KYC and onboarding benchmark research.",
        "Forrester. Onboarding and dispute-handling benchmarks.",
        "Aite-Novarica. Commercial onboarding TAT industry data.",
    ]),
    ("GOVERNANCE & REGULATION", [
        "Federal Reserve. SR 11-7: Guidance on Model Risk Management.",
        "OCC. Heightened Standards for Large Banks.",
        "European Union. Artificial Intelligence Act, in force 2024.",
        "UK FCA. AI principles and discussion papers.",
    ]),
]
top = BODY_TOP + Inches(0.9)
col_w = (BODY_W - Inches(0.6)) / 2
col_h = (BODY_H - Inches(1.0)) / 2
for idx, (head, items) in enumerate(sources):
    col = idx % 2
    row = idx // 2
    x = MARGIN + col * (col_w + Inches(0.6))
    y = top + row * (col_h + Inches(0.2))
    # Mono header
    text(s, x, y, col_w, Inches(0.3), head,
         font=FONT_MONO, size=9, bold=True, color="BLUE",
         letter_spacing=0.22)
    rect(s, x, y + Inches(0.32), Inches(0.4), Inches(0.025), "BLUE")
    # Items
    iy = y + Inches(0.5)
    for it in items:
        text(s, x, iy, col_w, Inches(0.5),
             "·  " + it,
             font=FONT_BODY, size=11, color="SLATE",
             line_spacing=1.45)
        iy += Inches(0.4)

footer(s, 25)


# ─────────────────────────────────────────────────────────────────────
# Save
#
# Repo layout (after the docs/ + src/ split):
#   src/decks/<slug>/build.py     ← this file
#   docs/decks/<slug>/<file>.pptx ← what we write (served by GitHub Pages)
#
# The build always writes into docs/ so the published site has the latest
# PPTX without a copy step.
# ─────────────────────────────────────────────────────────────────────
import os
THIS = os.path.dirname(os.path.abspath(__file__))
out_dir = os.path.abspath(
    os.path.join(THIS, "..", "..", "..", "docs", "decks",
                 "agentic-ai-banking-pov"))
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, "agentic-ai-banking-pov.pptx")
prs.save(out)
print(f"Saved {out}")
print(f"Slides: {len(prs.slides)}")
